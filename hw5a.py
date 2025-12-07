import streamlit as st
import pdfplumber
from transformers import pipeline
from pptx import Presentation
from PIL import Image
import numpy as np
import streamlit.components.v1 as components
from html import escape

st.set_page_config(page_title="AI Generated Detector", layout="centered")

@st.cache_resource
def load_detector(model_name: str = "roberta-base-openai-detector"):
    return pipeline(
        "text-classification",
        model=model_name,
        truncation=True,
    )

@st.cache_resource
def get_ocr():
    return load_ocr()

ocr = get_ocr()

def detect_ai(text: str, detector):
    """Return aggregated label, aggregated score, and per-chunk info list.

    per_chunk: list of dicts {'chunk': str, 'label': str, 'score': float}
    """
    if not text or not text.strip():
        return "UNKNOWN", 0.0, []
    chunk_size = 512
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    per_chunk = []
    for chunk in chunks:
        try:
            r = detector(chunk)[0]
            per_chunk.append({
                "chunk": chunk,
                "label": r.get("label", "UNKNOWN"),
                "score": float(r.get("score", 0.0)),
            })
        except Exception:
            per_chunk.append({"chunk": chunk, "label": "ERROR", "score": 0.0})
    if not per_chunk:
        return "UNKNOWN", 0.0, []
    scores = {}
    for p in per_chunk:
        lbl = p["label"]
        scores.setdefault(lbl, []).append(p["score"])
    avg_scores = {lbl: sum(v) / len(v) for lbl, v in scores.items()}
    best_label = max(avg_scores, key=avg_scores.get)
    return best_label, avg_scores[best_label], per_chunk

def extract_ppt_text(file):
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def extract_pdf_text(file):
    texts = []
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n".join(texts)

def extract_image_text(file):
    img = Image.open(file).convert("RGB")
    result = ocr.ocr(np.array(img))
    texts = []
    for page in result:
        for line in page:
            try:
                texts.append(line[1][0])
            except Exception:
                continue
    return "\n".join(texts)

st.title("🧠 AI 偵測系統（文字 / PPT / PDF / 圖片）")

# Sidebar: model and display options
st.sidebar.header("設定")
model_options = [
    "roberta-base-openai-detector",
    "roberta-large-openai-detector",
    "custom",
]
model_choice = st.sidebar.selectbox("選擇檢測模型", model_options, index=0)
custom_model = ""
if model_choice == "custom":
    custom_model = st.sidebar.text_input("自訂模型名稱 (Hugging Face)", "")
model_name = custom_model.strip() if model_choice == "custom" and custom_model.strip() else model_choice

show_chunk_texts = st.sidebar.checkbox("顯示 chunk 原文 (展開式)", value=False)
max_chunks_show = st.sidebar.slider("最大顯示 chunk 數量", min_value=1, max_value=50, value=10)

# load detector for chosen model
try:
    detector = load_detector(model_name)
except Exception as e:
    st.sidebar.error(f"載入模型失敗：{e}")
    detector = load_detector("roberta-base-openai-detector")

mode = st.selectbox("選擇輸入類型", ["文字", "PPTX", "PDF", "圖片"]) 

content = ""

if mode == "文字":
    content = st.text_area("輸入文字", height=250)

elif mode == "PPTX":
    file = st.file_uploader("上傳 PPTX", type=["pptx"])
    if file:
        try:
            content = extract_ppt_text(file)
        except Exception as e:
            st.error(f"無法解析 PPTX：{e}")

elif mode == "PDF":
    file = st.file_uploader("上傳 PDF", type=["pdf"])
    if file:
        try:
            content = extract_pdf_text(file)
        except Exception as e:
            st.error(f"無法解析 PDF：{e}")

elif mode == "圖片":
    file = st.file_uploader("上傳圖片", type=["png", "jpg", "jpeg"])
    if file:
        try:
            content = extract_image_text(file)
            st.image(file, caption="上傳圖片")
        except Exception as e:
            st.error(f"無法解析圖片：{e}")

if st.button("開始偵測"):
    if not content or not content.strip():
        st.warning("沒有可用的文字來源，請輸入文字或上傳檔案。")
    else:
        label, score, per_chunk = detect_ai(content, detector)
        st.subheader("📊 偵測結果")
        st.metric("判定", label)
        try:
            st.progress(int(score * 100))
        except Exception:
            st.progress(0)
        st.write(f"信心分數：**{score:.2%}**")

        with st.expander(f"每個 chunk 的判定細節 (顯示前 {max_chunks_show} 個)"):
            search = st.text_input("在 chunk 中搜尋文字或標籤", value="")
            if search and search.strip():
                s = search.lower()
                filtered = [p for p in per_chunk if s in p['chunk'].lower() or s in p['label'].lower()]
            else:
                filtered = per_chunk

            display_chunks = filtered[:max_chunks_show]
            for idx, p in enumerate(display_chunks, start=1):
                cols = st.columns([1, 4, 1])
                cols[0].write(f"**{idx}**")
                cols[1].write(f"Label = **{p['label']}**, Score = **{p['score']:.2%}**")
                try:
                    cols[0].progress(int(p['score'] * 100))
                except Exception:
                    cols[0].progress(0)

                # Copy button using an HTML component so client clipboard can be used
                safe_text = escape(p['chunk'])
                html = f"""
                <div style='display:flex; gap:8px; align-items:center;'>
                  <button onclick="navigator.clipboard.writeText(document.getElementById('ta_{idx}').value)">Copy</button>
                  <textarea id='ta_{idx}' style='display:none;'>{safe_text}</textarea>
                </div>
                """
                components.html(html, height=40)

                if show_chunk_texts:
                    with st.expander(f"顯示 Chunk {idx} 原文"):
                        st.code(p['chunk'])

            if len(filtered) > max_chunks_show:
                st.write(f"... 總共 {len(filtered)} 個符合的 chunk，僅顯示前 {max_chunks_show} 個。")
